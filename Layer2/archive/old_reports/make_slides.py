"""
Build an ECG safety pipeline summary deck (PowerPoint).
Run from the ECG Processing root directory.

Output: Results/slides/ECG_Safety_Pipeline.pptx
"""
from __future__ import annotations

import io
from pathlib import Path
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
RES      = ROOT / "Results"
OUT_DIR  = RES / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_OUT = OUT_DIR / "ECG_Safety_Pipeline.pptx"

FI_DIR   = RES / "feature_importance"

# ── brand colours ────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x23, 0x4E)   # slide background accents
C_BLUE   = RGBColor(0x27, 0x80, 0xB9)   # Layer 2
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)   # abnormal / warning
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # healthy / pass
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # caution
C_GRAY   = RGBColor(0xEC, 0xF0, 0xF1)   # slide background
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x2C, 0x3E, 0x50)


# ── helpers ──────────────────────────────────────────────────────────────────
W   = Inches(13.33)   # 16:9 slide width
H   = Inches(7.5)     # 16:9 slide height
ML  = Inches(0.5)     # margin left
MT  = Inches(1.1)     # margin top (below title bar)
CW  = Inches(12.33)   # content width


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # totally blank
    return prs.slides.add_slide(layout)


def bg_rect(slide, color: RGBColor, l=0, t=0, w=None, h=None):
    w = w or W; h = h or H
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def title_bar(slide, title: str, subtitle: str = ""):
    bg_rect(slide, C_NAVY, 0, 0, W, Inches(0.9))
    tx = slide.shapes.add_textbox(ML, Inches(0.08), W - ML * 2, Inches(0.75))
    tf = tx.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.bold = True; r.font.size = Pt(24); r.font.color.rgb = C_WHITE
    if subtitle:
        r2 = p.add_run(); r2.text = "   " + subtitle
        r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)


def text_box(slide, text: str, l, t, w, h,
             size=Pt(11), bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    return tx


def colored_rect(slide, color, l, t, w, h, text="", tsize=Pt(10), tcol=C_WHITE):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    if text:
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = tsize; r.font.bold = True; r.font.color.rgb = tcol
    return sh


def fig_to_pptx(slide, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    slide.shapes.add_picture(buf, l, t, width=w, height=h)
    plt.close(fig)


def img_to_pptx(slide, path: Path, l, t, w, h):
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)


def section_divider(prs: Presentation, title: str, subtitle: str = "") -> None:
    sl = blank_slide(prs)
    bg_rect(sl, C_NAVY)
    tx = sl.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(1.2))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = C_WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)


# ── metric card ──────────────────────────────────────────────────────────────
def metric_card(slide, label: str, value: str, l, t, w=Inches(2.2), h=Inches(1.1),
                color=C_BLUE):
    colored_rect(slide, color, l, t, w, h)
    text_box(slide, value, l, t + Inches(0.12), w, Inches(0.55),
             size=Pt(26), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text_box(slide, label, l, t + Inches(0.62), w, Inches(0.42),
             size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)


# ── table helper ─────────────────────────────────────────────────────────────
def add_table(slide, headers, rows, l, t, w, h,
              header_color=C_NAVY, row_colors=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    col_w = w // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    # Header
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True; run.font.size = Pt(10)
                run.font.color.rgb = C_WHITE

    # Rows
    alt = [RGBColor(0xF4, 0xF6, 0xF7), RGBColor(0xFF, 0xFF, 0xFF)]
    for i, row in enumerate(rows):
        rc = row_colors[i] if row_colors else alt[i % 2]
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rc
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = C_DARK


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_NAVY)
    # Diagonal accent
    sh = sl.shapes.add_shape(6, Inches(9), Inches(0), Inches(5), H)
    sh.fill.solid(); sh.fill.fore_color.rgb = C_BLUE
    sh.line.fill.background()

    text_box(sl, "ECG Safety Pipeline", Inches(0.7), Inches(1.8), Inches(8), Inches(1.4),
             size=Pt(42), bold=True, color=C_WHITE)
    text_box(sl, "Layer 1 · Layer 2 · Cross-Dataset Validation",
             Inches(0.7), Inches(3.1), Inches(8), Inches(0.6),
             size=Pt(18), color=RGBColor(0xBD, 0xC3, 0xC7))
    text_box(sl, "Closed-loop cardiomyoplasty / MNA stimulation safety",
             Inches(0.7), Inches(3.75), Inches(8), Inches(0.5),
             size=Pt(13), color=RGBColor(0x85, 0x92, 0x9E))
    text_box(sl, "Master Thesis  ·  June 2026",
             Inches(0.7), Inches(6.6), Inches(4), Inches(0.5),
             size=Pt(11), color=RGBColor(0x85, 0x92, 0x9E))


def slide_overview(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Safety Architecture", "Three-layer pipeline")

    layers = [
        ("Layer 1", "Fast R-peak detection\n& RR supervisor",
         "Causal filtering → candidate peaks\nRefractory / blanking rules\nImmediate inhibit on gross failures",
         C_RED),
        ("Layer 2", "Handcrafted feature\nsafety gate",
         "40 features: timing, morphology,\nrhythm variability, wavelet energy\nPer-session healthy baseline calibration",
         C_BLUE),
        ("Layer 3", "SSL / Anomaly\ndetection (research)",
         "Self-supervised encoder\nDeep SVDD / Mahalanobis in embedding\nOffline / server-side only",
         C_ORANGE),
    ]
    for i, (title, sub, desc, col) in enumerate(layers):
        x = ML + i * Inches(4.3)
        colored_rect(sl, col, x, MT, Inches(4.0), Inches(1.1), text=title, tsize=Pt(16))
        text_box(sl, sub,  x, MT + Inches(1.15), Inches(4.0), Inches(0.7),
                 size=Pt(12), bold=True, color=col)
        text_box(sl, desc, x, MT + Inches(1.85), Inches(4.0), Inches(1.1),
                 size=Pt(10), color=C_DARK)

    # Decision rule
    text_box(sl, "Decision rule:",
             ML, Inches(4.1), Inches(2), Inches(0.4), size=Pt(11), bold=True, color=C_DARK)
    colored_rect(sl, C_NAVY, ML, Inches(4.5), CW, Inches(0.55),
                 text="permit  =  Layer1_safe  AND  Layer2_safe  AND  Layer3_safe_or_unavailable",
                 tsize=Pt(12))
    text_box(sl,
             "Safety rule: uncertainty = inhibit  |  false permit is the primary safety metric  |  Layer 3 optional",
             ML, Inches(5.15), CW, Inches(0.4), size=Pt(10), color=C_DARK)

    # Focus note
    colored_rect(sl, RGBColor(0xFD, 0xF2, 0xE9), ML, Inches(5.65), CW, Inches(0.55),
                 text="This deck focuses on Layer 2 validation (Layers 1 & 3 provide context)",
                 tsize=Pt(11), tcol=C_ORANGE)


def slide_layer2_features(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Layer 2 — Feature Set", "40 features across 3 clinical groups")

    groups = [
        ("Timing (RR)", C_RED,
         ["Beat coupling ratio  ★★★",
          "Mean / HR / SDNN / RMSSD",
          "RR CV, range, min/max",
          "Short / long RR fraction",
          "Beat-to-beat diff (mean, max)"]),
        ("Morphology", C_BLUE,
         ["Template correlation  ★★",
          "Neighbor beat correlation  ★★",
          "QRS width  ★★",
          "Beat amplitude",
          "Amplitude vs local median",
          "Post/pre-QRS area ratio"]),
        ("Signal & Wavelet", C_GREEN,
         ["RMS, peak-to-peak, energy",
          "Line length, zero-cross rate",
          "HF / LF noise ratio (SQI)",
          "Wavelet D1–D4 log energy",
          "Wavelet A log energy",
          "Wavelet D1–A entropy"]),
    ]
    for i, (grp, col, feats) in enumerate(groups):
        x = ML + i * Inches(4.3)
        colored_rect(sl, col, x, MT, Inches(4.0), Inches(0.55), text=grp, tsize=Pt(13))
        for j, f in enumerate(feats):
            text_box(sl, ("  " if not f.startswith("•") else "") + "• " + f,
                     x, MT + Inches(0.6 + j * 0.42), Inches(4.0), Inches(0.40),
                     size=Pt(10), color=C_DARK)

    text_box(sl, "★★★ = AUROC 0.95  ★★ = AUROC 0.80–0.93  (healthy vs ventricular beats, INCART)",
             ML, Inches(6.8), CW, Inches(0.4), size=Pt(10),
             color=RGBColor(0x64, 0x74, 0x7B), bold=False)


def slide_layer2_gate(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Layer 2 — Gate Logic", "Decision pipeline per beat")

    steps = [
        ("1  Hard rules", C_RED,
         "Immediate veto if any feature breaks a\nphysiological bound (e.g. coupling < 0.80,\ntemplate_corr < 0.55, SQI > threshold)"),
        ("2  Max Z-score", C_ORANGE,
         "Inhibit if any feature deviates strongly\nfrom healthy baseline\n(frozen quantile 0.90 from calibration)"),
        ("3  Mahalanobis\n    distance", C_BLUE,
         "Inhibit if joint feature vector is too far\nfrom healthy calibration covariance\n(threshold from healthy p99.9 quantile)"),
        ("4  PERMIT", C_GREEN,
         "All gates passed —\nstimulation permitted"),
    ]
    for i, (step, col, desc) in enumerate(steps):
        x = ML + i * Inches(3.15)
        colored_rect(sl, col, x, MT, Inches(2.9), Inches(1.0), text=step, tsize=Pt(13))
        text_box(sl, desc, x, MT + Inches(1.05), Inches(2.9), Inches(1.0),
                 size=Pt(10), color=C_DARK)
        if i < 3:
            text_box(sl, "→", x + Inches(2.95), MT + Inches(0.3), Inches(0.25), Inches(0.5),
                     size=Pt(18), bold=True, color=C_DARK)

    # Calibration note
    colored_rect(sl, C_NAVY, ML, Inches(3.3), CW, Inches(0.85),
                 text="Calibration: first 60% of record's healthy beats\n"
                      "→ estimate per-patient mean / covariance / thresholds   (NO abnormal labels used)",
                 tsize=Pt(11))

    # Hard rules table
    text_box(sl, "Key hard rules:", ML, Inches(4.3), Inches(3), Inches(0.4),
             size=Pt(11), bold=True, color=C_DARK)
    rules = [
        ("Beat coupling ratio",    "< 0.80",  "Too-early beat (PVC signature)"),
        ("Template correlation",   "< 0.55",  "QRS shape too different from template"),
        ("RR short fraction",      "> 50%",   "Rhythm dominated by short intervals"),
        ("HF noise ratio (raw)",   "> 15%",   "Signal quality gate (SQI)"),
    ]
    add_table(sl, ["Feature", "Threshold", "Meaning"],
              [[r[0], r[1], r[2]] for r in rules],
              ML, Inches(4.75), CW, Inches(1.65))


def slide_feature_importance(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Feature Importance", "AUROC: healthy vs ventricular beats — INCART oracle")

    img = FI_DIR / "feature_importance.png"
    if img.exists():
        img_to_pptx(sl, img, ML, MT, Inches(7.5), Inches(5.6))

    # Top-5 summary panel on the right
    x_r = Inches(8.3)
    text_box(sl, "Top 5 features", x_r, MT, Inches(4.7), Inches(0.4),
             size=Pt(12), bold=True, color=C_DARK)
    top5 = [
        ("Beat coupling ratio",    "0.948", C_RED),
        ("Neighbor beat corr.",    "0.932", C_BLUE),
        ("Template correlation",   "0.898", C_BLUE),
        ("QRS width",              "0.801", C_BLUE),
        ("RR coeff. of variation", "0.791", C_ORANGE),
    ]
    for j, (name, auc, col) in enumerate(top5):
        y = MT + Inches(0.5 + j * 0.78)
        colored_rect(sl, col, x_r, y, Inches(4.7), Inches(0.65))
        text_box(sl, name, x_r + Inches(0.1), y + Inches(0.05),
                 Inches(3.3), Inches(0.35), size=Pt(10), color=C_WHITE)
        text_box(sl, f"AUROC {auc}", x_r + Inches(3.3), y + Inches(0.05),
                 Inches(1.3), Inches(0.35), size=Pt(11), bold=True,
                 color=C_WHITE, align=PP_ALIGN.RIGHT)

    text_box(sl,
             "Timing features (coupling) are the strongest single discriminator.\n"
             "Morphology adds specificity. Rhythm variability captures episode context.",
             x_r, Inches(5.8), Inches(4.7), Inches(0.8),
             size=Pt(10), color=C_DARK)


def slide_pca(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "PCA of Layer 2 Features", "Healthy vs ventricular beats — INCART oracle")

    pca_img   = FI_DIR / "pca_2d.png"
    load_img  = FI_DIR / "pca_loadings.png"
    if pca_img.exists():
        img_to_pptx(sl, pca_img, ML, MT, Inches(6.0), Inches(5.2))
    if load_img.exists():
        img_to_pptx(sl, load_img, Inches(6.7), MT, Inches(6.3), Inches(5.2))

    text_box(sl,
             "PC1 separates healthy from ventricular beats.\n"
             "Main drivers: coupling ratio, QRS width, template correlation.\n"
             "PC2 captures rhythm variability (SDNN, RMSSD).",
             ML, Inches(6.3), CW, Inches(0.55), size=Pt(10), color=C_DARK)


def slide_mitbih_results(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "MIT-BIH Results (dev set)", "Oracle mode — zero-shot gate")

    metric_card(sl, "Healthy Permit",   "84.3%", ML,               MT, color=C_GREEN)
    metric_card(sl, "Abnormal Inhibit", "91.5%", ML + Inches(2.5), MT, color=C_RED)
    metric_card(sl, "False Permit",      "8.5%", ML + Inches(5.0), MT, color=C_ORANGE)
    metric_card(sl, "SVT Inhibit",      "35.3%", ML + Inches(7.5), MT, color=C_BLUE)
    metric_card(sl, "Records",           "48",   ML + Inches(10.0), MT, color=C_NAVY)

    # Inhibit breakdown note
    text_box(sl, "Inhibit class breakdown (abnormal beats):", ML, Inches(2.7), CW, Inches(0.35),
             size=Pt(11), bold=True, color=C_DARK)
    breakdown = [
        ("Hard rules (coupling, template, SQI)", "24.4%", "Immediate veto — no calibration needed"),
        ("Z-score gate",                         "~40%",  "Strong feature deviation from healthy baseline"),
        ("Mahalanobis gate",                     "~12%",  "Joint distribution distance"),
        ("False permits (within_baseline)",       "8.5%",  "Subtle beats that resemble healthy pattern"),
    ]
    add_table(sl, ["Inhibit class", "Rate", "Interpretation"],
              breakdown, ML, Inches(3.1), CW, Inches(1.6))

    text_box(sl,
             "Oracle gap vs adaptive L1: ~4.3 pp  |  Gap concentrated in 10 bad Layer 1 records\n"
             "→ Layer 2 is not the bottleneck; R-peak detection limits the adaptive result",
             ML, Inches(5.0), CW, Inches(0.55), size=Pt(10), color=C_DARK)

    text_box(sl, "Calibration: healthy-only (60% of record)  |  No abnormal labels used",
             ML, Inches(5.65), CW, Inches(0.4), size=Pt(10), bold=True, color=C_BLUE)


def _ds_row(dataset, role, hp, ai, fp, svt):
    return [dataset, role, f"{hp:.1%}", f"{ai:.1%}", f"{fp:.1%}", f"{svt:.1%}"]


def slide_cross_dataset_baseline(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Cross-Dataset Baseline", "Zero-shot oracle — frozen MIT-BIH gate, healthy-only calibration")

    rows = [
        _ds_row("MIT-BIH",  "Dev set",           0.8435, 0.9149, 0.0851, 0.3535),
        _ds_row("INCART",   "Morphology shift",  0.8785, 0.8167, 0.1833, 0.1613),
        _ds_row("NSTDB",    "Noise stress",      0.7826, 0.8094, 0.1906, 0.1736),
        _ds_row("SVDB",     "Irregular rhythm",  0.8782, 0.7635, 0.2365, 0.3405),
    ]
    row_colors = [
        RGBColor(0xEA, 0xF4, 0xEB),
        RGBColor(0xFD, 0xF2, 0xE9),
        RGBColor(0xFD, 0xF2, 0xE9),
        RGBColor(0xFD, 0xF2, 0xE9),
    ]
    add_table(sl, ["Dataset", "Role", "HP", "AI", "FP", "SVT inhibit"],
              rows, ML, MT, CW, Inches(1.8), row_colors=row_colors)

    text_box(sl, "Key observations:", ML, Inches(3.35), CW, Inches(0.35),
             size=Pt(11), bold=True, color=C_DARK)

    obs = [
        ("Healthy permit", C_GREEN,
         "Remains 78–88% across all datasets — gate does not over-inhibit normal beats externally."),
        ("Abnormal inhibit", C_RED,
         "Drops on INCART (81.7%) and SVDB (76.3%). Root cause: subtle short-coupled PVCs with near-normal morphology."),
        ("False permit", C_ORANGE,
         "Rises to 18–24% on external sets. 98% of false permits are V (PVC) beats — missed due to coupling rule dropout."),
        ("NSTDB noise", C_BLUE,
         "Performance reliable at SNR >= 12 dB (AI ~100%). Below 6 dB: SQI gate should pre-screen."),
    ]
    for i, (label, col, desc) in enumerate(obs):
        y = Inches(3.8 + i * 0.7)
        colored_rect(sl, col, ML, y, Inches(2.1), Inches(0.55), text=label, tsize=Pt(10))
        text_box(sl, desc, ML + Inches(2.2), y + Inches(0.08), Inches(10.1), Inches(0.45),
                 size=Pt(10), color=C_DARK)


def slide_incart_root_cause(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "INCART Root-Cause Analysis", "Why 18.3% false permit on morphology-shift dataset?")

    # Left: root causes
    causes = [
        ("1. Subtle PVCs", C_RED,
         "98% of false permits are V beats with\ntemplate_corr ≈ 0.93  (look like normal QRS)\nbut coupling ≈ 0.54  (arrive too early)"),
        ("2. Coupling rule silently dropped", C_ORANGE,
         "If rr__beat_coupling_ratio is NaN-heavy in\nthe 60% calibration window, the hard rule\nis removed from the calibrator — PVCs escape"),
        ("3. Mahalanobis weak on INCART", C_BLUE,
         "AUROC 0.57 (near random) — per-patient\ncalibration makes same-patient PVCs look\nin-distribution in covariance space"),
    ]
    for i, (title, col, desc) in enumerate(causes):
        y = MT + i * Inches(1.55)
        colored_rect(sl, col, ML, y, Inches(2.0), Inches(1.2), text=title, tsize=Pt(12))
        text_box(sl, desc, ML + Inches(2.1), y + Inches(0.1), Inches(4.2), Inches(1.0),
                 size=Pt(10), color=C_DARK)

    # Right: score table
    text_box(sl, "Feature AUROC on INCART (healthy vs ventricular):",
             Inches(7.0), MT, Inches(5.8), Inches(0.4), size=Pt(11), bold=True, color=C_DARK)
    auroc_rows = [
        ("Beat coupling ratio",  "0.948", "★★★"),
        ("Neighbor beat corr",   "0.932", "★★★"),
        ("Template correlation", "0.898", "★★"),
        ("QRS width",            "0.801", "★★"),
        ("Mahalanobis distance", "0.571", "weak"),
    ]
    add_table(sl, ["Feature", "AUROC", ""],
              auroc_rows, Inches(7.0), MT + Inches(0.45), Inches(5.8), Inches(1.8))

    text_box(sl, "Fix: enforce coupling rule even when NaN in cal window → coupling becomes a non-droppable gate",
             ML, Inches(5.85), CW, Inches(0.45), size=Pt(10), bold=True, color=C_RED)

    colored_rect(sl, C_NAVY, ML, Inches(6.35), CW, Inches(0.45),
                 text="Global z-score tightening costs ~22–41 pp HP  |  Coupling rule costs only ~1–2 pp HP",
                 tsize=Pt(11))


def slide_coupling_sweep(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Coupling Threshold Sweep", "0.75 / 0.80 / 0.85 — all datasets, oracle mode")

    # Table
    rows_data = [
        ("MIT-BIH",  "84.3%", "91.5%", "94.6%", "94.8%", "95.0%"),
        ("INCART",   "87.8%", "81.7%", "94.6%", "95.8%", "96.8%"),
        ("NSTDB",    "78.3%", "80.9%", "99.4%", "99.7%", "99.9%"),
        ("SVDB",     "87.8%", "76.3%", "90.4%", "92.9%", "94.5%"),
    ]
    add_table(sl,
              ["Dataset", "HP (all)", "AI — baseline", "AI — 0.75", "AI — 0.80", "AI — 0.85"],
              rows_data, ML, MT, CW, Inches(1.9))

    text_box(sl, "HP change at coupling 0.80 (vs baseline):",
             ML, Inches(3.55), Inches(4.5), Inches(0.4), size=Pt(11), bold=True, color=C_DARK)
    hp_rows = [
        ("MIT-BIH", "84.3%", "82.2%", "−2.1 pp"),
        ("INCART",  "87.8%", "86.6%", "−1.2 pp"),
        ("NSTDB",   "78.3%", "78.3%", "0 pp"),
        ("SVDB",    "87.8%", "86.5%", "−1.3 pp"),
    ]
    add_table(sl, ["Dataset", "Baseline HP", "HP @ 0.80", "Change"],
              hp_rows, ML, Inches(3.95), Inches(5.5), Inches(1.7))

    text_box(sl, "SVT / premature beat inhibit rises sharply (35–93% depending on dataset).",
             Inches(6.3), Inches(3.55), Inches(6.7), Inches(0.4),
             size=Pt(11), bold=True, color=C_ORANGE)
    text_box(sl,
             "Coupling rule is not ventricular-specific:\n"
             "short-coupled SVT/PAC beats also trigger it.\n\n"
             "This is a deliberate conservative safety policy:\n"
             "any premature beat → do not stimulate.",
             Inches(6.3), Inches(4.05), Inches(6.7), Inches(1.6),
             size=Pt(10), color=C_DARK)


def slide_final_results(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Final Results — Coupling 0.80", "Oracle mode · zero-shot · all datasets")

    # 4 metric cards per dataset
    datasets = [
        ("MIT-BIH",  0.822, 0.948, 0.052, 0.569),
        ("INCART",   0.866, 0.958, 0.042, 0.933),
        ("NSTDB",    0.783, 0.997, 0.003, 0.807),
        ("SVDB",     0.865, 0.929, 0.071, 0.824),
    ]
    for i, (ds, hp, ai, fp, svt) in enumerate(datasets):
        y = MT + i * Inches(1.4)
        text_box(sl, ds, ML, y + Inches(0.35), Inches(1.4), Inches(0.5),
                 size=Pt(13), bold=True, color=C_DARK)
        metric_card(sl, "HP",  f"{hp:.1%}", ML + Inches(1.5), y, Inches(2.4), Inches(1.1), C_GREEN)
        metric_card(sl, "AI",  f"{ai:.1%}", ML + Inches(4.1), y, Inches(2.4), Inches(1.1), C_RED)
        metric_card(sl, "FP",  f"{fp:.1%}", ML + Inches(6.7), y, Inches(2.4), Inches(1.1), C_ORANGE)
        metric_card(sl, "SVT", f"{svt:.1%}", ML + Inches(9.3), y, Inches(2.4), Inches(1.1), C_BLUE)

    colored_rect(sl, C_NAVY, ML, Inches(7.0), CW, Inches(0.3),
                 text="HP = Healthy permit   AI = Abnormal inhibit   FP = False permit   SVT = Premature non-ventricular inhibit",
                 tsize=Pt(9))


def slide_nstdb_snr(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "NSTDB Noise Stress", "Performance by SNR level — oracle mode")

    snr_rows = [
        ("+24 dB", "76.3%", "100%",  "0.0%",  "Safe"),
        ("+18 dB", "76.0%", "100%",  "0.0%",  "Safe"),
        ("+12 dB", "79.8%", "99.8%", "0.2%",  "Safe"),
        ("+6 dB",  "84.2%", "73.7%", "26.3%", "Below threshold"),
        ("0 dB",   "83.1%", "65.9%", "34.1%", "Below threshold"),
        ("−6 dB",  "70.2%", "46.3%", "53.7%", "Below threshold"),
    ]
    row_colors = (
        [RGBColor(0xEA, 0xF4, 0xEB)] * 3 +
        [RGBColor(0xFD, 0xEB, 0xEB)] * 3
    )
    add_table(sl, ["SNR", "HP", "AI", "FP", "Status"],
              snr_rows, ML, MT, Inches(8.0), Inches(2.5), row_colors=row_colors)

    text_box(sl, "Operational envelope: SNR >= 12 dB",
             ML, Inches(4.1), Inches(6), Inches(0.45),
             size=Pt(14), bold=True, color=C_GREEN)
    text_box(sl,
             "Primary NSTDB noise type is electrode motion (low-frequency, 1–10 Hz).\n"
             "This corrupts RR features and bypass the HF-SQI gate — caught instead by max_zscore.\n"
             "Below SNR +6 dB an upstream SQI pre-check should blanket-inhibit stimulation.",
             ML, Inches(4.65), CW, Inches(0.9), size=Pt(10), color=C_DARK)

    # Visual SNR bar
    snr_ai = [1.0, 1.0, 0.998, 0.737, 0.659, 0.463]
    snr_labels = ["+24", "+18", "+12", "+6", "0", "-6"]
    fig, ax = plt.subplots(figsize=(4.5, 2.0))
    colors = ["#27ae60"] * 3 + ["#e74c3c"] * 3
    ax.bar(snr_labels, snr_ai, color=colors, edgecolor="white")
    ax.axhline(0.90, color="orange", linestyle="--", linewidth=1.5, label="90% target")
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Abnormal inhibit", fontsize=8)
    ax.set_xlabel("SNR (dB)", fontsize=8)
    ax.set_title("AI by SNR (coupling 0.80)", fontsize=9, fontweight="bold")
    ax.legend(fontsize=7)
    ax.tick_params(labelsize=7)
    plt.tight_layout()
    fig_to_pptx(sl, fig, Inches(9.2), Inches(3.8), Inches(3.8), Inches(2.8))


def slide_conclusions(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_NAVY)
    text_box(sl, "Key Conclusions", ML, Inches(0.35), W - ML * 2, Inches(0.6),
             size=Pt(28), bold=True, color=C_WHITE)

    points = [
        (C_GREEN,  "Healthy permit is robust across datasets (78–87%) after per-session calibration."),
        (C_RED,    "Abnormal inhibit is strong on MIT-BIH (91.5%) and after coupling fix reaches >= 92% everywhere."),
        (C_ORANGE, "The coupling rule is the single most powerful gate: AUROC 0.95, costs only ~1–2 pp HP."),
        (C_BLUE,   "Three complementary feature groups: timing, morphology, rhythm variability — each adds value."),
        (C_RED,    "INCART false permits: subtle PVCs with near-normal morphology but short coupling (~0.54).\n    Fix: keep coupling rule active at all times (do not drop on NaN-heavy cal window)."),
        (C_ORANGE, "SVT / PAC beats are heavily inhibited by coupling rule — a deliberate conservative policy.\n    Must be stated explicitly in any deployment claim."),
        (C_BLUE,   "Operational limits: SNR >= 12 dB for reliable performance; upstream SQI gate below that."),
        (C_GREEN,  "All results are zero-shot (no external abnormal labels). This is a valid generalization claim."),
    ]
    for i, (col, text) in enumerate(points):
        y = Inches(1.1 + i * 0.74)
        sh = sl.shapes.add_shape(1, ML, y + Inches(0.15), Inches(0.25), Inches(0.25))
        sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()
        text_box(sl, text, ML + Inches(0.4), y, Inches(12.5), Inches(0.7),
                 size=Pt(11), color=C_WHITE)


def slide_next_steps(prs):
    sl = blank_slide(prs)
    bg_rect(sl, C_GRAY)
    title_bar(sl, "Next Steps", "Immediate priorities")

    items = [
        ("1", C_RED,    "Fix coupling hard-rule persistence",
         "Prevent rr__beat_coupling_ratio from being dropped when NaN-heavy in calibration window.\nEnsure it is always evaluated as a hard gate at inference time."),
        ("2", C_ORANGE, "Re-run full benchmark with coupling 0.80 enforced",
         "Confirm 92–99% AI on MIT-BIH / INCART / NSTDB / SVDB with correct gate.\nReport both oracle and adaptive Layer 1 results."),
        ("3", C_BLUE,   "Rat ECG baseline calibration workflow",
         "Design per-session calibration protocol: 2–5 min healthy baseline recording.\nValidate that per-rat feature distributions are stable before arrhythmia induction."),
        ("4", C_GREEN,  "Supervisor discussion",
         "Present cross-dataset results. Agree on false-permit budget for animal experiments.\nConfirm whether SVT-class inhibition is acceptable in the target animal model."),
    ]
    for i, (num, col, title, desc) in enumerate(items):
        y = MT + i * Inches(1.3)
        colored_rect(sl, col, ML, y, Inches(0.5), Inches(1.1), text=num, tsize=Pt(20))
        text_box(sl, title, ML + Inches(0.65), y, CW - Inches(0.65), Inches(0.45),
                 size=Pt(12), bold=True, color=col)
        text_box(sl, desc, ML + Inches(0.65), y + Inches(0.45), CW - Inches(0.65), Inches(0.65),
                 size=Pt(10), color=C_DARK)


# ═════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs);                       print("  1/11  Title")
    slide_overview(prs);                    print("  2/11  Architecture overview")
    slide_layer2_features(prs);             print("  3/11  Layer 2 features")
    slide_layer2_gate(prs);                 print("  4/11  Gate logic")
    slide_feature_importance(prs);          print("  5/11  Feature importance")
    slide_pca(prs);                         print("  6/11  PCA")
    slide_mitbih_results(prs);              print("  7/11  MIT-BIH results")
    slide_cross_dataset_baseline(prs);      print("  8/11  Cross-dataset baseline")
    slide_incart_root_cause(prs);           print("  9/11  INCART root cause")
    slide_coupling_sweep(prs);              print(" 10/11  Coupling sweep")
    slide_final_results(prs);              print(" 11/11  Final results")
    slide_nstdb_snr(prs);                   print(" 12/13  NSTDB SNR")
    slide_conclusions(prs);                 print(" 13/13  Conclusions")
    slide_next_steps(prs);                  print(" 14/14  Next steps")

    prs.save(str(PPTX_OUT))
    print(f"\nSaved: {PPTX_OUT}")


if __name__ == "__main__":
    main()

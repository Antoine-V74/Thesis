"""
Improve Midterm_improved_EPFL_MNA.pptx:
  - ground the narrative in literature (citation footers + a few body edits)
  - correct the lookahead result numbers to the full uncapped cross-dataset runs
  - add slides: 'What I tried', 'NSTDB noise robustness', 'References'
  - renumber page-number placeholders
Preserves the EPFL template styling by editing in place and cloning styled slides.
"""
import copy
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path(r"C:\Users\antoi\OneDrive\Bureau\Master Thesis")
PPTX = BASE / "Power Point" / "Midterm_improved_EPFL_MNA.pptx"
SLIDES_DIR = Path(r"C:\Users\antoi\OneDrive\Bureau\Master Thesis\Code Base\ECG Processing\Results\slides")
CHART_PNG = SLIDES_DIR / "lookahead_results_full.png"

EPFL_RED = "#FF0000"
DARK = "#2B2B2B"
MIDGRAY = "#9A9A9A"

# ---------------------------------------------------------------------------
# 1. Accurate lookahead results figure
#    fast_causal_gated, all features, mean over the common datasets
#    (mitdb, nstdb, svdb) that were run at every lookahead.
# ---------------------------------------------------------------------------
LOOKAHEADS = ["+50 ms", "+100 ms", "+150 ms"]
HP = [29.6, 69.0, 71.0]   # healthy permit
AI = [96.5, 96.7, 95.8]   # abnormal inhibit
FP = [3.5, 3.3, 4.2]      # false permit


def make_chart():
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    x = np.arange(len(LOOKAHEADS))
    w = 0.26
    fig, ax = plt.subplots(figsize=(10.6, 3.7), dpi=150)
    b1 = ax.bar(x - w, HP, w, label="Healthy permit (want high)", color=EPFL_RED)
    b2 = ax.bar(x, AI, w, label="Abnormal inhibit (want high)", color=DARK)
    b3 = ax.bar(x + w, FP, w, label="False permit (want low)", color=MIDGRAY)
    for bars in (b1, b2, b3):
        for r in bars:
            ax.annotate(f"{r.get_height():.0f}%",
                        (r.get_x() + r.get_width() / 2, r.get_height()),
                        textcoords="offset points", xytext=(0, 3),
                        ha="center", va="bottom", fontsize=10, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(LOOKAHEADS, fontsize=13, fontweight="bold")
    ax.set_ylim(0, 110)
    ax.set_ylabel("Rate (%)", fontsize=11)
    ax.set_yticks([0, 25, 50, 75, 100])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.legend(loc="lower center", ncol=3, frameon=False, fontsize=10,
              bbox_to_anchor=(0.5, -0.28))
    fig.tight_layout()
    fig.savefig(CHART_PNG, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("Saved chart:", CHART_PNG)


# ---------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------
def set_para_text(para, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.text = text


def replace_on_slide(slide, mapping):
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            key = para.text.strip()
            if key in mapping:
                set_para_text(para, mapping[key])


def add_footer(slide, text, top=6.62, size=9):
    tb = slide.shapes.add_textbox(Inches(0.62), Inches(top), Inches(11.9), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.italic = True
    f.color.rgb = RGBColor(0x80, 0x80, 0x80)
    return tb


def duplicate_slide(prs, source):
    # `source` may be an index or a slide object
    if isinstance(source, int):
        source = prs.slides[source]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)
    for shp in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def move_slide(prs, slide, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    # find the sldId that points to this slide
    rid = None
    for sid in ids:
        if sid.get(qn("r:id")) == _rid_of(prs, slide):
            rid = sid
            break
    if rid is None:
        return
    lst.remove(rid)
    lst.insert(new_index, rid)


from pptx.oxml.ns import qn


def _rid_of(prs, slide):
    # map slide part to its rId in the presentation part
    part = slide.part
    for rid, rel in prs.part.rels.items():
        if rel.target_part is part:
            return rid
    return None


def renumber(prs):
    for i, slide in enumerate(prs.slides):
        for ph in slide.placeholders:
            if ph.has_text_frame and ph.text_frame.text.strip().isdigit():
                set_para_text(ph.text_frame.paragraphs[0], str(i + 1))


def keep_only_header(slide, header_text):
    """Delete body shapes; keep placeholders, the header band textbox and its line."""
    for shp in list(slide.shapes):
        st = shp.shape_type
        txt = shp.text_frame.text.strip() if shp.has_text_frame else ""
        is_placeholder = (st == 14 or st is None and shp.is_placeholder if hasattr(shp, "is_placeholder") else st == 14)
        keep = False
        if st == 14:
            keep = True
        elif txt == header_text:
            keep = True
        elif st == 9:  # LINE
            try:
                if abs(shp.top / 914400 - 1.94) < 0.15:
                    keep = True
            except Exception:
                keep = False
        if not keep:
            shp._element.getparent().remove(shp._element)


def add_text(slide, left, top, width, height, lines, size=12, bold_first=False,
             color="#2B2B2B", align=None):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.color.rgb = RGBColor.from_string(color.lstrip("#"))
        if bold_first and i == 0:
            f.bold = True
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
    return tb


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------
def main():
    make_chart()
    prs = Presentation(str(PPTX))
    S = prs.slides

    # ---- Slide 2 outline: fix +200 mention -------------------------------
    replace_on_slide(S[1], {
        "4. Results: +50/+100/+150/+200 ms": "4. What I tried + results (+50/+100/+150 ms)",
    })

    # ---- Slide 5 (idx4): cite Song 2025 in the body ----------------------
    replace_on_slide(S[4], {
        "MNA force is reported as roughly 2-3x more fatigue resistant than conventional skeletal-muscle grafts.":
            "Song et al. (2025) report a myoneural actuator with about 260% greater fatigue resistance and closed-loop control demonstrated in a rodent model.",
    })

    # ---- Slide 14 (idx13): tighten interpretation to match full runs -----
    replace_on_slide(S[13], {
        "Too early for reliable morphology: healthy permit collapses, so the gate becomes overly conservative.":
            "Too early for reliable morphology: healthy permit collapses to ~30%, so the gate becomes overly conservative (abnormal inhibit and false permit stay good).",
        "Best main operating point: enough post-R ECG to score morphology while keeping decision latency acceptable.":
            "Best operating point: healthy permit ~69%, abnormal inhibit ~97%, false permit ~3%. Enough post-R ECG to score morphology at acceptable latency.",
        "Similar safety metrics, but little added value unless hardware timing or mechanical delay requires the extra margin.":
            "Marginal: healthy permit ~71% but false permit creeps up (~4%) and latency grows. +200 ms (capped sweep) is similar. Little added value.",
    })

    # ---- Citation footers on the narrative + results + next-steps slides --
    footers = {
        2:  "Refs: Roche et al. 2017 (Sci Transl Med); Bonnemain et al. 2022 (Annu Rev Biomed Eng) - direct cardiac compression supports the heart without contacting blood.",
        3:  "Refs: Geddes et al. 1993; Lucas et al. 1992; Mario et al. 1999; Salmons & Jarvis 2009 - cardiomyoplasty was limited by a coupled actuator / timing / relaxation problem, not one cause.",
        4:  "Ref: Song et al. 2025 - fatigue-resistant myoneural actuator (~260% fatigue resistance; closed-loop rodent control).",
        5:  "Design: hierarchical control = fast trigger / beat-to-beat / supervisory loops (project ECG-gating spec; cf. Magkoutas et al. 2022, adaptive assist control).",
        6:  "Principle: safety-first layering inspired by lightweight upstream anomaly filtering; ML may veto, never command stimulation.",
        9:  "Method: personalised healthy-baseline gate - robust z-score + Ledoit-Wolf Mahalanobis distance on handcrafted ECG features.",
        13: "Refs: Geddes et al. 1993 - subject-specific optimal R-delay ~58 ms (40-80 ms); cardiac electromechanical window ~60-100 ms, consistent with +100 ms being the first robust point.",
        19: "Refs: Lucas et al. 1992 (relaxation); Li et al. (muscle peak after LV peak); Geddes et al. 1993 (subject-specific optimum, wrap-tightness confound).",
    }
    for idx, txt in footers.items():
        add_footer(S[idx], txt)

    # capture slide references before any cloning/insertion shifts indices
    nextsteps_slide = S[14]   # next-steps template (also gets a footer later)
    tmpl_noise = S[11]        # Layer-1 results (stat-box look)
    tmpl_refs = S[16]         # Questions slide (clean header band)

    # ---- Slide 13 (idx12): rebuild body with accurate chart --------------
    s13 = S[12]
    keep_only_header(s13, "RESULTS - LOOKAHEAD COMPARISON")
    s13.shapes.add_picture(str(CHART_PNG), Inches(0.75), Inches(2.15),
                           width=Inches(8.1))
    add_text(s13, 9.15, 2.30, 3.1, 0.4, "Result choice", size=13, bold_first=True,
             color="#FF0000")
    add_text(s13, 9.15, 2.85, 3.1, 3.0,
             ["+100 ms is the main deployment point: healthy permit recovers to ~69% "
              "while abnormal inhibit stays ~97% and false permit ~3%.",
              "",
              "+50 ms is too early for morphology; +150/+200 ms add latency with little gain."],
             size=11)
    add_footer(s13, "fast_causal_gated, all features. +50/+100/+150 ms = mean over the common "
                    "full-run datasets (mitdb, nstdb, svdb). Source: cross_dataset_causal_*ms/overall_summary.csv.",
               top=5.75)

    # ---- Slide 25 (idx24): correct the numeric table ---------------------
    s25 = S[24]
    table = None
    for shp in s25.shapes:
        if shp.has_table:
            table = shp.table
            break
    if table is not None:
        rows = [
            ("+50 ms",  "29.6%", "96.5%", "3.5%", "Too early; healthy permit collapses"),
            ("+100 ms", "69.0%", "96.7%", "3.3%", "Main deployment point"),
            ("+150 ms", "71.0%", "95.8%", "4.2%", "Marginal gain, more latency"),
            ("+200 ms", "75.0%", "93.7%", "6.3%", "Similar (capped sweep only)"),
        ]
        for r_i, vals in enumerate(rows, start=1):
            for c_i, val in enumerate(vals):
                cell = table.cell(r_i, c_i)
                para = cell.text_frame.paragraphs[0]
                set_para_text(para, val)
    replace_on_slide(s25, {
        "Mode: fast_causal_gated, all features, quick multi-dataset lookahead comparison. Values are aggregated from the existing Layer 2 result CSVs.":
            "Mode: fast_causal_gated, all features. +50/+100/+150 ms from full uncapped cross-dataset runs (mitdb, nstdb, svdb, mean). +200 ms from an earlier capped sweep.",
    })

    # ---- NEW slide: 'What I tried' (clone next-steps numbered template) ---
    tried = duplicate_slide(prs, nextsteps_slide)   # clone next-steps look
    replace_on_slide(tried, {
        "NEXT STEPS": "WHAT I TRIED",
        "Next step: connect ECG safety to mechanical benefit":
            "What I compared before choosing +100 ms",
        "Pressure data": "Windowing",
        "Bring arterial/LV pressure into the TDT pipeline and define beat-level mechanical labels.":
            "Centered +/-2.5 s oracle was an offline, non-causal upper bound. Moved to a causal [R-w, R+lookahead] window for deployment.",
        "Animal baseline": "Decision timing",
        "Calibrate thresholds per animal/session before stimulation or arrhythmia induction.":
            "Compared gated (full L2 this beat), next-beat permit, and stateful risk carry-over. Gated is cleanest; the others trade latency.",
        "Hardware timing": "Feature set",
        "Measure total delay: detection, computation, pulse generator, nerve-muscle activation.":
            "All features vs signal-only. All-features keeps abnormal inhibit high; signal-only mainly rejects noise and misses morphology.",
        "Closed-loop test": "Thresholds",
        "Compare fixed vs adaptive stimulation delay using pressure/flow response.":
            "Swept the Mahalanobis scale and the RR coupling threshold (operating curves) to push false permit down.",
        "Goal for the next phase: prove that the safe trigger also lands in the mechanically useful window.":
            "Bottom line: +100 ms, causal, gated, all-features, with a tuned coupling/Mahalanobis threshold is the most defensible operating point.",
    })
    add_footer(tried, "Modes & windows from run_cross_dataset_benchmark.py; thresholds from the operating-curve sweep (per_beat.csv / coupling_sweep).")
    move_slide(prs, tried, 14)   # place right after slide 14 (interpretation)

    # add the next-steps citation footer AFTER cloning so the clone never inherits it
    add_footer(nextsteps_slide,
               "Refs: pressure labelling Geddes et al. 1993; PEP error budget Richer et al. 2025 (PEPbench); "
               "impedance surrogate Osswald et al. 2000; adaptive control Magkoutas et al. 2022.")

    # ---- NEW slide: NSTDB noise robustness (clone stat-box template 12) ---
    noise = duplicate_slide(prs, tmpl_noise)   # clone Layer-1 results look (4 stat boxes)
    replace_on_slide(noise, {
        "RESULTS - R-PEAK DETECTION": "RESULTS - NOISE ROBUSTNESS (NSTDB)",
        "Layer 1 supports causal validation":
            "Layer 2 stays safe down to clinically usable noise",
        "99.9%": "100%",
        "NSRDB sensitivity": "abnormal inhibit, SNR >= 12 dB",
        "98.9%": "0%",
        "NSRDB PPV": "false permit, SNR >= 12 dB",
        "100%": "97%",
        "NSTDB SNR>=12 sensitivity": "abnormal inhibit, SNR +6 dB",
        "2.8 ms": "~83%",
        "confirmation delay": "abnormal inhibit, SNR <= 0 dB",
        "Interpretation":
            "Interpretation",
        "The fast causal detector localizes R peaks with high sensitivity. Remaining issues are mainly duplicate/noisy candidates and recovery behavior, which the RR supervisor and Layer 2 are designed to handle.":
            "At clinically usable SNR (>= +12 dB) the gate inhibits 100% of abnormal beats with zero false permit. Performance only degrades under severe noise (<= 0 dB).",
        "Why this is enough for the midterm":
            "Safety design implication",
        "The project can now evaluate safety and timing with a deployable trigger source instead of an offline oracle.":
            "Under severe noise, false permit rises, so the system must lean on Layer-1 SQI inhibition: uncertainty defaults to inhibit.",
        "Sources: fast_causal_threshold_three_dataset_summary.csv; layer1_detector_comparison_with_adaptive_gated_nstdb_snr_ge12.csv":
            "Source: cross_dataset_causal_100ms/nstdb_snr_breakdown.csv (fast_causal_gated, all features).",
    })
    move_slide(prs, noise, 26)   # supplementary, after the numeric table area

    # ---- NEW slide: References (clone the 'Questions' slide) --------------
    refs = duplicate_slide(prs, tmpl_refs)   # clone Questions slide (clean header band)
    replace_on_slide(refs, {
        "QUESTIONS": "REFERENCES",
        "Thank you": "Key references",
    })
    # remove the two body textboxes of the cloned questions slide
    for shp in list(refs.shapes):
        if shp.has_text_frame:
            t = shp.text_frame.text.strip()
            if t.startswith("Supplementary slides follow") or t.startswith("They are organized"):
                shp._element.getparent().remove(shp._element)
    ref_lines = [
        "Geddes et al. (1993). Importance of timing muscle contraction in dynamic cardiomyoplasty. PACE.",
        "Lucas et al. (1992). Importance of muscle relaxation in dynamic cardiomyoplasty. PACE.",
        "Li et al. Optimization of programming in dynamic cardiomyoplasty (muscle/ventricle synchronization).",
        "Mario et al. (1999). Demand dynamic cardiomyoplasty: two-year results.",
        "Carpentier et al. (1993). Dynamic cardiomyoplasty at seven years. J Thorac Cardiovasc Surg.",
        "Salmons & Jarvis (2009). Cardiomyoplasty reviewed: lessons from the past, prospects for the future.",
        "Elefteriades (2022). A brief history of cardiomyoplasty: worth another look? Rev Cardiovasc Med.",
        "Patel et al. Dynamic cardiomyoplasty: insights into the mechanisms of its success.",
        "Yilmaz et al. (2003). Dynamic cardiomyoplasty: impact of effective pacing. Int J Cardiol.",
        "Roche et al. (2017). Soft robotic sleeve supports heart function. Sci Transl Med.",
        "Horvath et al. (2018). Coupling of a soft robotic sleeve to the heart. Ann Biomed Eng.",
        "Bonnemain et al. (2022). Direct cardiac compression devices. Annu Rev Biomed Eng.",
        "Song et al. (2025). A fatigue-resistant myoneural actuator for implantable biohybrid systems.",
        "Osswald et al. (2000). Intracardiac impedance and RV dP/dtmax (closed-loop stimulation sensor).",
        "Richer et al. (2025). PEPbench: benchmarking automated pre-ejection-period extraction. Psychophysiology.",
        "Magkoutas et al. (2022). Physiologic data-driven iterative learning control for LVADs. Front Cardiovasc Med.",
    ]
    add_text(refs, 0.9, 2.15, 11.5, 5.0, ref_lines, size=11)
    # refs is appended last already; leave it at the end

    renumber(prs)
    out = PPTX
    try:
        prs.save(str(PPTX))
    except PermissionError:
        out = PPTX.with_name("Midterm_improved_EPFL_MNA_v2.pptx")
        prs.save(str(out))
        print("Original locked (open in PowerPoint). Saved to:", out)
    print("Saved deck:", out, "| slides:", len(prs.slides))


if __name__ == "__main__":
    main()

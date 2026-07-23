#!/usr/bin/env python3
"""Generate a supervisor-facing Layer 3 slide deck (PPTX + Markdown outline)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPORTS = Path(__file__).resolve().parents[1] / "reports"
PPTX_OUT = REPORTS / "LAYER3_SUPERVISOR_DECK.pptx"
MD_OUT = REPORTS / "LAYER3_SUPERVISOR_DECK.md"


def _add_slide(prs: Presentation, title: str, bullets: list[str], notes: str = "") -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(17)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Layer 3 — Learned ECG Anomaly Safety Veto"
    slide.placeholders[1].text = (
        "Master Thesis — ECG-triggered cardiac stimulation\n"
        "Supervisor briefing | July 2026"
    )
    slide.notes_slide.notes_text_frame.text = (
        "Layer 3 is the optional research extension of the three-layer safety pipeline. "
        "It never commands stimulation; it only permits or inhibits. "
        "Human PhysioNet validation is proxy validation before rat deployment."
    )

    _add_slide(
        prs,
        "Project context: three-layer safety pipeline",
        [
            "Layer 1 — fast deterministic R-peak detection + RR supervisor (always on)",
            "Layer 2 — handcrafted ECG feature gate vs healthy baseline (primary deploy path)",
            "Layer 3 — learned embedding anomaly veto (optional / server-side research)",
            "Final rule: permit only if L1 AND L2 AND (L3 if enabled) all permit",
            "Uncertainty, failure, or missing data → inhibit (fail-safe)",
        ],
        "Layer 2 is the likely first rat-deployment gate. Layer 3 asks whether SSL "
        "representations can add morphology sensitivity beyond handcrafted features.",
    )

    _add_slide(
        prs,
        "What Layer 3 does (and does not do)",
        [
            "Input: short beat-synchronous ECG window around a candidate R-peak",
            "Output: permit or inhibit — a safety veto only",
            "NOT a clinical arrhythmia classifier; does not diagnose VT/VF/AF",
            "NOT required for the immediate safety loop; server failure → inhibit",
            "Goal: detect when morphology deviates from a session healthy baseline",
        ],
        "The operational question is: does this beat look close enough to the "
        "calibrated healthy ECG to allow stimulation?",
    )

    _add_slide(
        prs,
        "Layer 3 runtime pipeline",
        [
            "ECG window → ECGEncoder1D (1D ResNet, ~700K params) → 128-d embedding",
            "Per-record healthy calibration → Mahalanobis or kNN distance score",
            "Threshold from held-out healthy scores (conformal default, α=10%)",
            "score ≤ threshold → permit | score > threshold → inhibit",
            "Optional: L2 normalize + PCA(32) + Ledoit-Wolf covariance shrinkage",
        ],
        "The scorer is transparent and cheap to re-fit per animal/session. "
        "Only the encoder weights transfer across subjects; the baseline is always personalized.",
    )

    _add_slide(
        prs,
        "Why beat-synchronous evaluation?",
        [
            "One trigger beat → one window → one embedding → one permit/inhibit",
            "Primary window: 8 s @ 125 Hz (rhythm + morphology; L3 must carry rhythm itself)",
            "Morphology ablation: 1 s window (local QRS) — secondary, not the primary claim",
            "Causal window + optional 50–150 ms post-R lookahead (stimulation latency sim)",
            "Oracle mode vs Layer 1 gated mode (upper bound vs closed-loop deployment)",
        ],
        "Fixed sliding windows dilute the current beat with surrounding rhythm. "
        "Beat-sync evaluation matches how the controller actually decides.",
    )

    _add_slide(
        prs,
        "Encoder pretraining (self-supervised + supervised)",
        [
            "ECGEncoder1D: 1D ResNet, GroupNorm, 128-d embedding",
            "SSL arms (A/A1/B/B1): unlabeled ECG + physiology-aware augmentations",
            "Arm C: supervised contrastive (SupCon) using public labels at PRETRAIN ONLY",
            "All heads (projection/expander/SupCon) discarded — deploy is label-free one-class",
            "Human pretrain → frozen encoder → per-record/session healthy baseline re-fit",
        ],
        "Key distinction: labels may shape the encoder offline (Arm C), but the deployed "
        "decision never sees a danger label — it only fits a healthy baseline per subject.",
    )

    _add_slide(
        prs,
        "Phase 1: encoder comparison arms (fixed downstream scorer)",
        [
            "A0 — Layer 2 handcrafted features → Mahalanobis/kNN (control floor)",
            "A — NT-Xent / CLOCS-style contrastive SSL (baseline learned encoder)",
            "A1 — VICReg non-contrastive SSL (no negatives, no reconstruction)",
            "B — MAE mask + same-window consistency (primary masked arm)",
            "B1 — MAE + subject/record contrastive (ablation; prefer healthy-only)",
            "C — Supervised contrastive (SupCon): public labels at pretrain, label-free deploy",
        ],
        "All arms share the identical scorer and threshold logic, so any difference is "
        "representation quality, not the decision rule. Multi-lead is a future appendix, not Arm C.",
    )

    _add_slide(
        prs,
        "Thresholding and safety metrics",
        [
            "Primary metric: false-permit rate on DANGEROUS beats (VT/VF/noise policy)",
            "Secondary: healthy-permit rate (therapy availability)",
            "Default threshold: conformal prediction (α=10% healthy false-inhibit budget)",
            "Legacy: healthy quantile threshold (e.g. 99th percentile)",
            "Wilson confidence intervals on DANGEROUS false-permit; coverage report CSV",
        ],
        "False permit = stimulation allowed during unsafe rhythm — the main safety error. "
        "Benign ectopy is reported separately (don't-care for false-permit accounting).",
    )

    _add_slide(
        prs,
        "Literature alignment: ZEROSHOT iEEG (Yu 2026)",
        [
            "Independent parallel design: SSL encoder → per-subject Mahalanobis on healthy baseline",
            "Their result: within-subject zero-label AUROC 0.954 (near supervised 0.964)",
            "Supports personalized baseline > population model for biosignal anomaly",
            "Our bar is higher: false-permit + CIs, not AUROC alone",
            "Anomaly ≠ danger: policy grouping and causal deployment matter for ECG safety",
        ],
        "The paper validates the architecture choice. Our novelty is the safety framing, "
        "policy-aware evaluation, and adaptation to ECG stimulation veto.",
    )

    _add_slide(
        prs,
        "Current implementation status",
        [
            "Encoder, augmentations, Mahalanobis/kNN scorer, conformal threshold — implemented",
            "Beat-sync + window validation harness with Phase 1 tables + record bootstrap — implemented",
            "SSL arms A (NT-Xent), A1 (VICReg), B (MAE+consistency), B1 (subject-contrastive) — implemented",
            "Arm C (SupCon supervised) — implemented; labels pretrain-only, encoder-only checkpoint",
            "Exploratory 8 s runs: SSL arms not yet clearly beating A0 → motivates Arm C + locked gold pilot",
        ],
        "End-to-end smoke test covers index build, pretrain (incl. supcon), validation, and comparison to A0. "
        "Headline safety numbers await the locked MIT-BIH gold pilot on the cluster.",
    )

    _add_slide(
        prs,
        "Validation plan (human → rat)",
        [
            "Human MIT-BIH / PhysioNet: labeled proxy validation with policy groups",
            "Per-record healthy calibration → score DANGEROUS / BENIGN / NORMAL separately",
            "Report worst-record and worst-fold behavior, not just means",
            "Rat deployment: 2–5 min healthy baseline → session calibration → live veto",
            "Prospective animal validation required before any therapy reliance on Layer 3",
        ],
        "Human labels select evaluation windows only. Deployable thresholds use healthy "
        "calibration exclusively — DANGEROUS labels are offline evaluation only.",
    )

    _add_slide(
        prs,
        "Next steps",
        [
            "Locked gold pilot: A0 + A + A1 + B (+B1 ablation) on MIT-BIH gold, matched config",
            "Add Arm C (SupCon) as the 'do labels help the representation?' test",
            "Phase 1 tables: false-permit on DANGEROUS at conformal α=10% + record-bootstrap CIs",
            "Compare beat-sync oracle vs Layer 1 gated trigger modes; CAV vs A0",
            "Rat baseline workflow once healthy rat ECG is available",
        ],
        "First deliverable is defensible safety tables comparing encoder families "
        "(incl. supervised Arm C), not AUROC. Multi-lead remains a future appendix, not Arm C.",
    )

    _add_slide(
        prs,
        "Questions for discussion",
        [
            "Can we access healthy rat ECG before main experiments?",
            "Single-lead or multi-lead rat acquisition?",
            "Anesthesia / condition for baseline ECG recording?",
            "Any induced arrhythmia segments available later for validation?",
            "Is an external Python server acceptable for optional Layer 3 analysis?",
            "Target false-permit budget for preliminary animal experiments?",
        ],
        "These questions determine calibration protocol, evaluation scope, "
        "and whether Layer 3 stays research-only or moves toward deployment.",
    )

    REPORTS.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))


def build_markdown() -> None:
    content = """# Layer 3 — Supervisor Slide Deck (outline)

Generated companion to `LAYER3_SUPERVISOR_DECK.pptx`.

---

## Slide 1 — Title

**Layer 3 — Learned ECG Anomaly Safety Veto**

Master Thesis — ECG-triggered cardiac stimulation  
Supervisor briefing | July 2026

---

## Slide 2 — Project context

- Layer 1: fast R-peak + RR supervisor (always on)
- Layer 2: handcrafted feature gate (primary deploy path)
- Layer 3: learned embedding anomaly veto (optional research)
- Final rule: permit only if L1 AND L2 AND (L3 if enabled)
- Uncertainty / failure → inhibit

---

## Slide 3 — What Layer 3 does

- Input: beat-synchronous ECG window
- Output: permit or inhibit (veto only)
- NOT a clinical classifier
- NOT required for immediate safety loop
- Detects deviation from session healthy baseline

---

## Slide 4 — Runtime pipeline

```
ECG → ECGEncoder1D → 128-d embedding → Mahalanobis/kNN → threshold → permit/inhibit
```

- ~700K param 1D ResNet
- Per-record healthy calibration
- Conformal threshold (default α=10%)
- Optional L2 + PCA(32) + Ledoit-Wolf

---

## Slide 5 — Beat-synchronous evaluation

- One beat → one decision
- Primary window: 8 s @ 125 Hz (rhythm + morphology)
- Morphology ablation: 1 s window (secondary)
- Causal + optional post-R lookahead (50–150 ms)
- Oracle vs Layer 1 gated modes

---

## Slide 6 — Encoder pretraining (SSL + supervised)

- SSL arms (A/A1/B/B1): unlabeled ECG + physiology-aware augmentations
- Arm C: supervised contrastive using public labels **at pretrain only**
- All training heads discarded → deploy is label-free one-class
- Human pretrain → frozen encoder → per-session healthy baseline re-fit

---

## Slide 7 — Phase 1 arms

| Arm | Method |
|-----|--------|
| A0 | Layer 2 handcrafted (control) |
| A | NT-Xent contrastive |
| A1 | VICReg non-contrastive |
| B | MAE + same-window consistency (primary masked) |
| B1 | MAE + subject-contrastive (ablation) |
| C | Supervised contrastive (SupCon) — labels pretrain-only |

Fixed downstream scorer isolates encoder quality. Multi-lead = future appendix, not Arm C.

---

## Slide 8 — Metrics

- **Primary:** false-permit on DANGEROUS (VT/VF/noise)
- **Secondary:** healthy-permit (availability)
- Conformal threshold + Wilson CIs
- Benign ectopy = don't-care for false-permit

---

## Slide 9 — ZEROSHOT literature

- Parallel architecture validated in iEEG (AUROC 0.954)
- Our bar: false-permit + CIs, not AUROC alone
- Personalized baseline > population model

---

## Slide 10 — Status

- Full pipeline implemented (encoder, scorer, validation, conformal threshold)
- SSL arms A / A1 / B / B1 implemented
- Arm C (SupCon supervised) implemented — labels pretrain-only, encoder-only checkpoint
- Exploratory 8 s runs: SSL not yet clearly beating A0 → motivates Arm C + locked pilot
- Smoke test covers index build, pretrain (incl. supcon), validation

---

## Slide 11 — Human → rat plan

- MIT-BIH proxy validation now
- Per-session rat calibration later
- Prospective animal validation required

---

## Slide 12 — Next steps

- Locked gold pilot: A0 / A / A1 / B (+B1 ablation) on MIT-BIH gold
- Add Arm C (SupCon): does using labels at pretrain help the representation?
- Phase 1 safety tables: false-permit + record-bootstrap CIs; CAV vs A0
- Rat baseline when data available

---

## Slide 13 — Questions

1. Healthy rat ECG access?
2. Single vs multi-lead?
3. Baseline recording conditions?
4. Induced arrhythmia data later?
5. External server acceptable?
6. False-permit budget for animals?

---

*See also: `reports/README.md`, `ALGORITHM_SUMMARY.md`, `LAYER3_ARCHITECTURE_RATIONALE.md`*
"""
    REPORTS.mkdir(parents=True, exist_ok=True)
    MD_OUT.write_text(content, encoding="utf-8")


def main() -> None:
    build_pptx()
    build_markdown()
    print(f"Wrote {PPTX_OUT}")
    print(f"Wrote {MD_OUT}")


if __name__ == "__main__":
    main()
